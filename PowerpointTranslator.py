from pptx import Presentation
import six
import os
import os.path
from os import path
from google.cloud import translate_v3 as translate
import sys
from progress.bar import Bar
import glob
import json

LOCATION = "us-central1"

GLOSSARY_ID = "my_first_glossary"


# Translate a file
def translate_file(filename):
    if not path.exists(filename):
        sys.exit("Source file for translation not found")

    prs = Presentation(filename)
    total_slide_cnt = len(prs.slides)
    print("Translating " + str(total_slide_cnt) + " slides in " + filename + " from " + str(source_lang_code) + " to " + str(target_lang_code))
    bar = Bar('Translating: ' + filename, max=total_slide_cnt)

    # cycle through each slide in the file
    for slide in prs.slides:
        for shape in slide.shapes:
            translate_shape(shape)
        bar.next()
    bar.finish()

    # save the translated file
    output_filename = filename.split('.')[0] + "_" + target_lang_code + "2.pptx"
    try:
        prs.save(output_filename)
    except:
        sys.exit("\n Couldn't save translated file, check if '" + output_filename + "' is already open in another application.")
    return


# Translate a string
def translate_text(text):
    # Note on calling translate_text method:
    #  1) Passing invalid target language codes will throw unexpected 403 missing a valid API key error.
    #  2) Translation API sometimes detects 'es' instead 'pt', passing source language improves translation quality.

    if isinstance(text, six.binary_type):
        text = text.decode("utf-8")

    output = ""
    if glossary_available:
        response = translate_client.translate_text(
            request={
                "parent": PARENT,
                "contents": [text],
                "mime_type": "text/plain",  # mime types: text/plain, text/html
                "source_language_code": source_lang_code,
                "target_language_code": target_lang_code,
                "glossary_config": glossary_config,
            }
        )

        for translation in response.glossary_translations:
            output = output + translation.translated_text
    else:
        response = translate_client.translate_text(
            request={
                "parent": PARENT,
                "contents": [text],
                "mime_type": "text/plain",  # mime types: text/plain, text/html
                "source_language_code": source_lang_code,
                "target_language_code": target_lang_code,
            }
        )

        # Display the translation for each input text provided
        for translation in response.translations:
            output = output + translation.translated_text

    return output


# Translate a paragraph
def translate_paragraph(paragraph):

    if (paragraph.text != ''):
        # Translate entire paragraph (typically a sentence), instead of translating a shorter phrase,
        # for a better quality translation.
        result = translate_text(paragraph.text)

        p = paragraph._p
        # remove all but the first run (instead of overwriting which could lose the font)
        for idx, run in enumerate(paragraph.runs):
            if idx == 0:
                continue
            p.remove(run._r)

        if (len(paragraph.runs) == 0):
            paragraph.add_run()

        # Entire paragraph gets font of the first run, if this was Wingdings then paragraph will be unreadable.
        # Typically this happens if paragraph is using a Wingdings special character such as a right-arrow,
        # with the remainder in a normal readable font. Therefore, check for this edge case and force Calibri instead.
        font = paragraph.runs[0].font
        if font.name == 'Wingdings':
            font.name = 'Calibri'

        paragraph.runs[0].text = result

    return


# Translate all text in a shape
def translate_shape(shape):
    if (shape.has_text_frame):
        # A shape which contains some text
        for paragraph in shape.text_frame.paragraphs:
            translate_paragraph(paragraph)

    elif(shape.has_table):
        # A Table
        for cell in shape.table.iter_cells():
            if not cell.is_spanned and cell.text != '':
                # In tables with merged cells, is_spanned checks if cell visible to user

                for paragraph in cell.text_frame.paragraphs:

                    translate_paragraph(paragraph)

    elif(shape.shape_type == 6):
        # Grouped shapes - use recursion
        for shp in shape.shapes:
            translate_shape(shp)

    return

# check arguments passed are valid
# assume pt to en translation unless overridden by arguments
source_lang_code = "pt"
target_lang_code = "en"
if len(sys.argv) == 4:
    source_lang_code = sys.argv[2].replace("-","")
    target_lang_code = sys.argv[3].replace("-","")

# fetch json key
my_GCP_key =  glob.glob("My_GCP_Translation_API_Key*.json" )
if len(my_GCP_key) != 1:
    sys.exit("Couldn't find key, file naming convention is 'My_GCP_Translation_API_Key*.json'. Exiting....")

json_file = open(my_GCP_key[0], )
json_data = json.load(json_file)
project_ID = json_data['project_id']
json_file.close()

os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = my_GCP_key[0]

PARENT = f"projects/{project_ID}/locations/{LOCATION}"

# create translation client and apply previously created glossary
translate_client = translate.TranslationServiceClient()
glossary = translate_client.glossary_path(project_ID, LOCATION, GLOSSARY_ID)
glossary_config = translate.TranslateTextGlossaryConfig(glossary=glossary)

# check if glossary includes the source and target languages
source_lang_code_glossary_available = False
target_lang_code_glossary_available = False
for glossary in translate_client.list_glossaries(parent=PARENT):
    if GLOSSARY_ID in glossary.name:
        # this is the glossary being used, check if it includes the source and target languages
        for language_code in glossary.language_codes_set.language_codes:
            if source_lang_code in language_code:
                source_lang_code_glossary_available = True
            if target_lang_code in language_code:
                target_lang_code_glossary_available = True

glossary_available = source_lang_code_glossary_available and target_lang_code_glossary_available
if glossary_available:
    print("Found a glossary for " + source_lang_code + " to " + target_lang_code)
else:
    print("No glossary found for " + source_lang_code + " to " + target_lang_code)

# make a list of pptx files to be translated
files_to_translate = []
if sys.argv[1] == "-all":
    # user wants to translate all pptx files in current folder
    all_pptx_files = glob.glob('*.pptx')
    # ignore any that are an output from a previous translation
    # eg if translating to en then anything ending in "_en2.pptx" is something we previously created
    ignore_1 = "_" + target_lang_code + ".pptx"
    ignore_2 = "_" + target_lang_code + "2.pptx"
    for file in all_pptx_files:
        if not(ignore_1 in file or ignore_2 in file):
            files_to_translate.append(file)
else:
    # user just wants to translate the named file
    files_to_translate.append(sys.argv[1])


for file in files_to_translate:
    translate_file(file)
